# app/tools/office_io.py

from __future__ import annotations

import os
import io
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urlparse
import mimetypes

import requests
import mammoth
from pptx import Presentation
from pptx.util import Inches, Pt

# Prefer PyMuPDF for PDF text (faster, better); fall back to pdfminer.
try:
    import fitz  # PyMuPDF
    _HAVE_FITZ = True
except Exception:
    _HAVE_FITZ = False
    from pdfminer.high_level import extract_text


# ---------- General Helpers ----------

def _run(cmd: List[str]) -> str:
    p = subprocess.run(cmd, capture_output=True, text=True)
    if p.returncode != 0:
        raise RuntimeError(f"{' '.join(cmd)}\n{p.stderr.strip()}")
    return p.stdout

def _guess_suffix(url: str, content_type: Optional[str]) -> str:
    if content_type:
        ext = mimetypes.guess_extension(content_type.split(";")[0].strip())
        if ext:
            return ext
    path = urlparse(url).path
    suf = Path(path).suffix
    return suf if suf else ""

def _ensure_dir(p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)


# ---------- Download / Inputs ----------

def fetch_to_tmp(inputs: Dict[str, Any]) -> str:
    """
    Accepts one of:
      - inputs["file_url"] or ["source_url"]
      - inputs["file_path"] (already local)
    Returns a local temp file path.
    """
    url = inputs.get("file_url") or inputs.get("source_url")
    if url:
        headers = {
            "User-Agent": os.getenv(
                "FETCH_USER_AGENT",
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/125.0.0.0 Safari/537.36"
            ),
            "Accept": "*/*",
        }
        try:
            with requests.get(url, headers=headers, stream=True, timeout=60, allow_redirects=True) as r:
                r.raise_for_status()
                suffix = _guess_suffix(url, r.headers.get("content-type"))
                fd, tmp_path = tempfile.mkstemp(suffix=suffix or "")
                with os.fdopen(fd, "wb") as f:
                    for chunk in r.iter_content(chunk_size=8192):
                        if chunk:
                            f.write(chunk)
                return tmp_path
        except requests.HTTPError as e:
            status = e.response.status_code if e.response is not None else "unknown"
            raise RuntimeError(f"Download failed ({status}) for {url}")
        except requests.RequestException as e:
            raise RuntimeError(f"Download error for {url}: {e}")

    file_path = inputs.get("file_path")
    if file_path and os.path.exists(file_path):
        tmp = Path(tempfile.mkdtemp()) / Path(file_path).name
        shutil.copy2(file_path, tmp)
        return str(tmp)

    raise ValueError("Provide file_url/source_url or file_path.")


# ---------- Outline Extraction ----------

def _text_to_outline(text: str, max_slides: int = 12) -> List[Dict[str, Any]]:
    """
    Very simple heuristic: split text into paragraphs and build slides with
    first line as title and following lines as bullets.
    """
    paras = [p.strip() for p in text.splitlines() if p.strip()]
    slides: List[Dict[str, Any]] = []
    chunk: List[str] = []
    for line in paras:
        # start new chunk at likely titles (short-ish lines or all-caps)
        if len(line) <= 80 or line.isupper():
            if chunk:
                title = chunk[0][:120]
                bullets = [s for s in chunk[1:6]]
                slides.append({"title": title, "bullets": bullets})
                if len(slides) >= max_slides:
                    break
                chunk = []
        chunk.append(line)

    if chunk and len(slides) < max_slides:
        title = chunk[0][:120]
        bullets = [s for s in chunk[1:6]]
        slides.append({"title": title, "bullets": bullets})

    # Fallback single slide if nothing parsed
    if not slides and paras:
        title = paras[0][:120]
        bullets = paras[1:6]
        slides = [{"title": title, "bullets": bullets}]

    return slides[:max_slides]

def pdf_to_outline(pdf_path: str, max_slides: int = 12) -> List[Dict[str, Any]]:
    if _HAVE_FITZ:
        doc = fitz.open(pdf_path)
        text_parts: List[str] = []
        for page in doc:
            text_parts.append(page.get_text("text"))
        text = "\n".join(text_parts)
    else:
        text = extract_text(pdf_path)
    return _text_to_outline(text, max_slides=max_slides)

def docx_to_outline(docx_path: str, max_slides: int = 12) -> List[Dict[str, Any]]:
    # Use mammoth to extract plain text from docx
    with open(docx_path, "rb") as f:
        result = mammoth.extract_raw_text(f)
    text = result.value or ""
    return _text_to_outline(text, max_slides=max_slides)

def pptx_to_outline(pptx_path: str, max_slides: int = 50) -> List[Dict[str, Any]]:
    prs = Presentation(pptx_path)
    outline: List[Dict[str, Any]] = []
    for slide in prs.slides:
        title = ""
        bullets: List[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = (shape.text or "").strip()
            if not txt:
                continue
            if not title:
                title = txt.splitlines()[0][:120]
            else:
                for line in txt.splitlines():
                    line = line.strip()
                    if line and line != title:
                        bullets.append(line)
        if title or bullets:
            outline.append({"title": title or "Slide", "bullets": bullets[:8]})
        if len(outline) >= max_slides:
            break
    return outline


# ---------- PPTX Builder ----------

def outline_to_pptx(outline: List[Dict[str, Any]], title: Optional[str] = None, dest_path: Optional[str] = None) -> str:
    prs = Presentation()
    # Title slide (optional)
    if title:
        layout = prs.slide_layouts[0]  # Title
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        try:
            slide.placeholders[1].text = ""
        except Exception:
            pass

    # Title & Content slides
    for s in outline:
        layout = prs.slide_layouts[1]  # Title & Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.get("title", "")[:120]
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = s.get("bullets") or []
        if bullets:
            body.text = str(bullets[0])[:300]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = str(b)[:300]
                p.level = 0

    if not dest_path:
        dest_path = tempfile.mktemp(suffix=".pptx")
    _ensure_dir(Path(dest_path))
    prs.save(dest_path)
    return dest_path


# ---------- LibreOffice Conversions ----------

def soffice_convert(src_path: str, fmt: str, outdir: Optional[str] = None) -> str:
    """
    fmt examples:
      - 'pdf'
      - 'html'
      - 'docx'
    Returns path to the converted file (or top-level html if html export).
    """
    outdir = outdir or tempfile.mkdtemp()
    cmd = ["soffice", "--headless", "--convert-to", fmt, "--outdir", outdir, src_path]
    _run(cmd)

    src = Path(src_path)
    if fmt == "html":
        # LibreOffice usually outputs {stem}.html next to assets directory
        out = Path(outdir) / f"{src.stem}.html"
        if out.exists():
            return str(out)
        # Some versions create an HTML directory — try to find first html
        for p in Path(outdir).glob("*.html"):
            return str(p)
        raise RuntimeError("HTML export did not produce an .html file")
    else:
        # map extension from fmt, crude but ok
        ext = "." + fmt.lower().split(":")[0]
        out = Path(outdir) / f"{src.stem}{ext}"
        if not out.exists():
            # Fallback: pick first matching file
            files = list(Path(outdir).glob(f"{src.stem}.*"))
            if files:
                out = files[0]
        if not out.exists():
            raise RuntimeError(f"Conversion to {fmt} failed (no output)")
        return str(out)

def zip_html_tree(html_path: str) -> str:
    root = Path(html_path).parent
    stem = Path(html_path).stem
    zip_path = Path(tempfile.mkdtemp()) / f"{stem}.zip"
    shutil.make_archive(str(zip_path.with_suffix("")), "zip", root)
    return str(zip_path)


# ---------- Artifact Saving ----------

def save_artifact(local_path: str, dest_name: Optional[str] = None) -> Tuple[str, str]:
    """
    Copies local_path into ARTIFACTS_DIR and returns:
      (absolute_dest_path, http_url)
    """
    try:
        from app.core.config import settings
        artifacts_dir = Path(getattr(settings, "ARTIFACTS_DIR", "/data/artifacts"))
        public_base = getattr(settings, "PUBLIC_BASE_URL", os.getenv("PUBLIC_BASE_URL", "http://localhost:8080"))
    except Exception:
        artifacts_dir = Path(os.getenv("ARTIFACTS_DIR", "/data/artifacts"))
        public_base = os.getenv("PUBLIC_BASE_URL", "http://localhost:8080")

    artifacts_dir.mkdir(parents=True, exist_ok=True)

    name = dest_name or Path(local_path).name
    dest = artifacts_dir / name
    # If the name exists, uniquify
    if dest.exists():
        dest = artifacts_dir / f"{Path(name).stem}_{next(tempfile._get_candidate_names())}{Path(name).suffix}"
    shutil.copy2(local_path, dest)

    url = f"{public_base.rstrip('/')}/artifacts/{dest.name}"
    return str(dest), url


# ---------- Convenience wrappers agents use ----------

def build_pptx_from_pdf(pdf_path: str, title: Optional[str], max_slides: int = 12) -> str:
    outline = pdf_to_outline(pdf_path, max_slides=max_slides)
    return outline_to_pptx(outline, title=title)

def build_pptx_from_docx(docx_path: str, title: Optional[str], max_slides: int = 12) -> str:
    outline = docx_to_outline(docx_path, max_slides=max_slides)
    return outline_to_pptx(outline, title=title)

def pptx_to_pdf_file(pptx_path: str) -> str:
    return soffice_convert(pptx_path, "pdf")

def pptx_to_docx_file(pptx_path: str) -> str:
    return soffice_convert(pptx_path, "docx")

def pptx_to_html5_zip(pptx_path: str) -> str:
    html = soffice_convert(pptx_path, "html")
    return zip_html_tree(html)


__all__ = [
    # fetch
    "fetch_to_tmp",
    # outline extractors
    "pdf_to_outline", "docx_to_outline", "pptx_to_outline",
    # build pptx
    "outline_to_pptx", "build_pptx_from_pdf", "build_pptx_from_docx",
    # conversions
    "pptx_to_pdf_file", "pptx_to_docx_file", "pptx_to_html5_zip",
    # artifacts
    "save_artifact",
]


# import os, shutil, subprocess, tempfile
# from pathlib import Path
# from typing import List, Dict, Any
# import urllib.request

# import mammoth
# from pdfminer.high_level import extract_text
# from pptx import Presentation

# def _run(cmd: List[str]):
#     p = subprocess.run(cmd, capture_output=True, text=True)
#     if p.returncode != 0:
#         raise RuntimeError(f"{' '.join(cmd)}\n{p.stderr}")
#     return p.stdout

# def fetch_to_tmp(inputs: Dict[str, Any]) -> str:
#     # Supports file_url | source_url | file_path
#     url = inputs.get("file_url") or inputs.get("source_url")
#     if url:
#         fd, path = tempfile.mkstemp()
#         os.close(fd)
#         urllib.request.urlretrieve(url, path)
#         return path
#     if inputs.get("file_path") and os.path.exists(inputs["file_path"]):
#         tmp = Path(tempfile.mkdtemp()) / Path(inputs["file_path"]).name
#         shutil.copy2(inputs["file_path"], tmp)
#         return str(tmp)
#     raise ValueError("Provide file_url/source_url or file_path.")

# def docx_to_outline(docx_path: str) -> List[Dict[str, Any]]:
#     with open(docx_path, "rb") as f:
#         html = mammoth.convert_to_html(f).value
#     sections = [s.strip() for s in html.split("<h") if "</h" in s]
#     outline = []
#     for s in sections:
#         title = s.split(">",1)[1].split("</h",1)[0].strip()
#         bullets = []
#         if "<li>" in s:
#             bullets = [x.split("</li>",1)[0] for x in s.split("<li>")[1:]]
#         outline.append({"title": title, "bullets": bullets[:6]})
#     return outline[:24]

# def pdf_to_outline(pdf_path: str) -> List[Dict[str, Any]]:
#     text = extract_text(pdf_path) or ""
#     chunks = [c.strip() for c in text.split("\n\n") if len(c.strip()) > 40]
#     return [{"title": c.split("\n")[0][:80], "bullets": c.split("\n")[1:6]} for c in chunks[:20]]

# def outline_to_pptx(outline: List[Dict[str, Any]], out_path: str, title: str = None):
#     prs = Presentation()
#     if title:
#         slide = prs.slides.add_slide(prs.slide_layouts[0])
#         slide.shapes.title.text = title
#         slide.placeholders[1].text = ""
#     for sec in outline:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text = sec.get("title","")[:120]
#         tf = slide.placeholders[1].text_frame
#         try:
#             tf.clear()
#         except Exception:
#             pass
#         for b in sec.get("bullets", [])[:12]:
#             p = tf.add_paragraph(); p.text = str(b)[:200]; p.level = 0
#     prs.save(out_path)
#     return out_path

# def soffice_convert(inp: str, fmt: str, outdir: str) -> str:
#     Path(outdir).mkdir(parents=True, exist_ok=True)
#     _run(["soffice","--headless","--convert-to",fmt,"--outdir",outdir,inp])
#     base = Path(inp).with_suffix(f".{fmt.split(':')[0]}")
#     return str(Path(outdir)/base.name)
